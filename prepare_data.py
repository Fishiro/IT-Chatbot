import os
import sys
import json
import time
import random
import hashlib
import tempfile
import subprocess

import docx
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document

# ============================================================
# CẤU HÌNH (có thể chỉnh qua biến môi trường, không cần sửa code)
# ============================================================
DATA_FOLDER = "learning_data"
INDEX_FOLDER = "faiss_index"
MANIFEST_PATH = "embed_manifest.json"

# Số đoạn văn bản gửi trong 1 lần gọi API embedding.
# Để nhỏ + có delay giữa các batch giúp KHÔNG bị vượt rate limit của Google.
EMBED_BATCH_SIZE = int(os.getenv("EMBED_BATCH_SIZE", "20"))
DELAY_BETWEEN_BATCHES = float(os.getenv("EMBED_BATCH_DELAY_SECONDS", "8"))

# Số lần thử lại tối đa khi gọi API embedding bị lỗi (429 / quota / mạng...)
MAX_RETRIES = int(os.getenv("EMBED_MAX_RETRIES", "6"))
RETRY_BASE_DELAY = float(os.getenv("EMBED_RETRY_BASE_DELAY", "5"))

LEGACY_EXTENSIONS = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}


def load_api_keys():
    """
    Đọc danh sách API key từ .env. Hỗ trợ 2 cách:
    - GEMINI_API_KEYS="key1,key2,key3,key4,key5"  (khuyên dùng, phân tách bằng dấu phẩy)
    - GEMINI_API_KEY="key1"                        (chỉ 1 key, giữ tương thích ngược)
    """
    multi = os.getenv("GEMINI_API_KEYS")
    if multi:
        keys = [k.strip() for k in multi.split(",") if k.strip()]
        if keys:
            return keys
    single = os.getenv("GEMINI_API_KEY")
    return [single] if single else []


class KeyRotator:
    """Xoay vòng nhiều API key: chủ động luân phiên qua các batch để chia đều tải,
    và khi 1 key bị rate-limit thì lập tức chuyển key khác thay vì ngồi chờ."""

    def __init__(self, api_keys):
        self.api_keys = api_keys
        self.models = [
            GoogleGenerativeAIEmbeddings(
                model="models/gemini-embedding-001", google_api_key=k)
            for k in api_keys
        ]
        self.n = len(self.models)
        self.idx = 0

    def current(self):
        return self.models[self.idx]

    def rotate(self):
        self.idx = (self.idx + 1) % self.n
        return self.current()


# ============================================================
# TIỆN ÍCH: CHUYỂN ĐỔI FILE OFFICE ĐỜI CŨ (.doc/.ppt/.xls)
# ============================================================
def convert_legacy_office_file(file_path, ext):
    """
    Chuyển .doc/.ppt/.xls sang .docx/.pptx/.xlsx để đọc được bằng
    python-docx / python-pptx / pandas.
    Ưu tiên LibreOffice (đa nền tảng, miễn phí) -> nếu không có thì
    thử Microsoft Office qua COM (chỉ chạy được trên Windows có cài Office).
    Trả về đường dẫn file đã convert, hoặc None nếu thất bại.
    """
    target_ext = LEGACY_EXTENSIONS[ext]
    tmpdir = tempfile.mkdtemp(prefix="legacy_convert_")
    errors = []

    # --- Cách 1: LibreOffice headless ---
    try:
        result = subprocess.run(
            ["soffice", "--headless", "--convert-to", target_ext.lstrip("."),
             "--outdir", tmpdir, file_path],
            timeout=120, capture_output=True, text=True,
            encoding="utf-8", errors="replace",
        )
        converted_path = os.path.join(
            tmpdir, os.path.splitext(os.path.basename(file_path))[
                0] + target_ext
        )
        if os.path.exists(converted_path):
            return converted_path
        errors.append(f"LibreOffice: exit code {result.returncode}, "
                      f"stderr: {result.stderr.strip()[:300]}")
    except FileNotFoundError:
        errors.append(
            "LibreOffice: không tìm thấy lệnh 'soffice' (chưa cài, hoặc chưa thêm vào PATH).")
    except subprocess.TimeoutExpired:
        errors.append("LibreOffice: quá thời gian chờ (timeout 120s).")
    except Exception as e:
        errors.append(f"LibreOffice: {e}")

    # --- Cách 2: Microsoft Office qua COM (chỉ Windows) ---
    try:
        import win32com.client as win32
        converted_path = os.path.join(
            tmpdir, os.path.splitext(os.path.basename(file_path))[
                0] + target_ext
        )
        abs_in = os.path.abspath(file_path)
        abs_out = os.path.abspath(converted_path)

        app = None
        try:
            if ext == ".doc":
                app = win32.Dispatch("Word.Application")
                app.Visible = False
                app.DisplayAlerts = False
                d = app.Documents.Open(abs_in, ReadOnly=True)
                d.SaveAs(abs_out, FileFormat=16)  # wdFormatXMLDocument (.docx)
                d.Close(False)
            elif ext == ".ppt":
                app = win32.Dispatch("PowerPoint.Application")
                app.DisplayAlerts = False
                p = app.Presentations.Open(abs_in, WithWindow=False)
                p.SaveAs(abs_out, 24)  # ppSaveAsOpenXMLPresentation
                p.Close()
            elif ext == ".xls":
                app = win32.Dispatch("Excel.Application")
                app.Visible = False
                app.DisplayAlerts = False
                wb = app.Workbooks.Open(abs_in, ReadOnly=True)
                wb.SaveAs(abs_out, FileFormat=51)  # xlOpenXMLWorkbook (.xlsx)
                wb.Close(False)
        finally:
            # Luôn đóng ứng dụng COM dù có lỗi/Ctrl+C giữa chừng,
            # tránh để lại tiến trình Word/PowerPoint/Excel chạy ngầm.
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass

        if os.path.exists(converted_path):
            return converted_path
        errors.append(
            "Microsoft Office COM: chạy xong nhưng không thấy file kết quả.")
    except ImportError:
        errors.append(
            "Microsoft Office COM: chưa cài pywin32 (pip install pywin32).")
    except Exception as e:
        errors.append(f"Microsoft Office COM: {e}")

    print(f"   [DEBUG] Chi tiết lỗi convert '{os.path.basename(file_path)}':")
    for err in errors:
        print(f"      - {err}")

    return None


# ============================================================
# HÀM HỖ TRỢ: ĐỌC ĐA ĐỊNH DẠNG TÀI LIỆU
# ============================================================
def ocr_scanned_pdf_pages(file_path, ocr_pages):
    """
    OCR các trang PDF không có text (dạng ảnh scan) bằng Tesseract.
    Cần cài Tesseract OCR (engine, không phải chỉ pip package) + gói ngôn ngữ tiếng Việt.
    Trả về chuỗi text OCR được, hoặc "" nếu không OCR được (thiếu Tesseract...).
    """
    try:
        import pytesseract
    except ImportError:
        print("   [!] Chưa cài pytesseract: pip install pytesseract")
        return ""

    # Cho phép chỉ định đường dẫn tesseract.exe qua .env nếu không nằm trong PATH
    tesseract_cmd = os.getenv("TESSERACT_CMD")
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

    try:
        pytesseract.get_tesseract_version()
    except Exception:
        print("   [!] Không tìm thấy Tesseract OCR engine trên máy. "
              "Cần cài riêng (không chỉ pip install) — xem hướng dẫn cài Tesseract OCR cho Windows.")
        return ""

    ocr_text = ""
    for page_num, page in ocr_pages:
        try:
            img = page.to_image(resolution=300).original
            page_ocr = pytesseract.image_to_string(img, lang="vie+eng")
            if page_ocr.strip():
                print(
                    f"      -> OCR trang {page_num}: nhận diện được {len(page_ocr.strip())} ký tự.")
                ocr_text += page_ocr + "\n"
        except Exception as e:
            print(f"      [!] OCR trang {page_num} lỗi: {e}")

    return ocr_text


def extract_text_from_file(file_path):
    """Tự động nhận diện đuôi file và trích xuất văn bản"""
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    try:
        # --- Office đời cũ: convert trước rồi đọc như bình thường ---
        if ext in LEGACY_EXTENSIONS:
            converted = convert_legacy_office_file(file_path, ext)
            if converted is None:
                print(f"   [!] Không thể chuyển đổi file cũ '{file_path}'. "
                      f"Cần cài LibreOffice (khuyên dùng) hoặc Microsoft Office + pywin32 (Windows), "
                      f"hoặc tự lưu lại (Save As) sang định dạng mới ({LEGACY_EXTENSIONS[ext]}).")
                return ""
            return extract_text_from_file(converted)

        # --- Word (.docx) ---
        if ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            # Đọc luôn nội dung trong bảng (table) nếu có
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(
                        cell.text for cell in row.cells if cell.text.strip())
                    if row_text.strip():
                        text += "\n" + row_text

        # --- PowerPoint (.pptx) ---
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"

        # --- Excel (.xlsx) ---
        elif ext == ".xlsx":
            xls = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in xls.items():
                text += f"\n--- Bảng: {sheet_name} ---\n"
                text += df.to_string(index=False) + "\n"

        # --- CSV ---
        elif ext == ".csv":
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)

        # --- Văn bản thuần ---
        elif ext == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        # --- PDF ---
        elif ext == ".pdf":
            import pdfplumber
            ocr_pages = []
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text += page_text + "\n"
                    else:
                        # Trang không có text -> có thể là ảnh scan, đánh dấu để OCR
                        ocr_pages.append((page_num, page))

            if ocr_pages:
                ocr_text = ocr_scanned_pdf_pages(file_path, ocr_pages)
                if ocr_text:
                    text += ocr_text
                elif not text.strip():
                    print(f"   [!] '{file_path}' là PDF dạng scan (ảnh) nhưng OCR thất bại hoặc chưa cài đặt. "
                          f"Xem hướng dẫn cài Tesseract OCR để đọc được loại file này.")

        else:
            print(f"   [!] Bỏ qua file định dạng không hỗ trợ: {file_path}")

    except Exception as e:
        print(f"   [X] Lỗi khi đọc file {file_path}: {e}")

    return text


# ============================================================
# TIỆN ÍCH: MANIFEST (theo dõi file nào đã nhúng, để chỉ xử lý phần MỚI)
# ============================================================
def load_manifest():
    if os.path.exists(MANIFEST_PATH):
        with open(MANIFEST_PATH, "r", encoding="utf-8") as f:
            return json.load(f)
    return {}


def save_manifest(manifest):
    with open(MANIFEST_PATH, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)


def compute_file_hash(path):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


# ============================================================
# TIỆN ÍCH: GỌI EMBEDDING API AN TOÀN (retry + backoff, tránh vượt quota)
# ============================================================
def embed_texts_with_retry(rotator, texts):
    for attempt in range(MAX_RETRIES):
        # Thử lần lượt qua TẤT CẢ key trong 1 "vòng" trước khi chịu ngồi chờ (sleep).
        for _ in range(rotator.n):
            model = rotator.current()
            key_no = rotator.idx + 1
            try:
                return model.embed_documents(texts)
            except Exception as e:
                msg = str(e).lower()
                is_rate_limit = any(
                    k in msg for k in ["429", "resource_exhausted", "quota", "rate limit", "resource exhausted"]
                )
                if is_rate_limit and rotator.n > 1:
                    print(
                        f"      [!] Key #{key_no}/{rotator.n} bị rate limit -> chuyển sang key tiếp theo...")
                    rotator.rotate()
                    continue  # thử ngay key kế tiếp, KHÔNG cần chờ
                elif is_rate_limit:
                    break  # chỉ có 1 key -> phải chờ backoff
                else:
                    print(
                        f"      [!] Lỗi khi gọi Embedding API (key #{key_no}): {e}")
                    # lỗi khác (không phải rate limit) -> chờ backoff rồi thử lại
                    break

        # Hết vòng key mà vẫn chưa xong -> chờ backoff rồi lặp lại vòng mới
        wait = RETRY_BASE_DELAY * (2 ** attempt) + random.uniform(0, 3)
        wait = min(wait, 120)
        print(f"      [!] Tất cả key vừa thử đều lỗi. Chờ {wait:.0f}s rồi thử lại "
              f"(vòng {attempt + 1}/{MAX_RETRIES})...")
        time.sleep(wait)
    raise RuntimeError(
        "Đã thử lại nhiều lần với toàn bộ API key nhưng vẫn thất bại. "
        "Tiến độ hiện tại đã được lưu (checkpoint) — chạy lại script sau để tiếp tục."
    )


def embed_documents_in_batches(rotator, texts):
    """Nhúng danh sách văn bản theo từng lô nhỏ, xoay vòng qua nhiều API key để tránh rate limit."""
    all_vectors = []
    total = len(texts)
    for i in range(0, total, EMBED_BATCH_SIZE):
        batch = texts[i:i + EMBED_BATCH_SIZE]
        if rotator.n > 1:
            rotator.rotate()  # chủ động đổi key mỗi batch để chia đều tải, không đợi lỗi mới đổi
        print(f"      -> Đang nhúng {i + 1}-{min(i + len(batch), total)}/{total} đoạn "
              f"(key #{rotator.idx + 1}/{rotator.n})...")
        vectors = embed_texts_with_retry(rotator, batch)
        all_vectors.extend(vectors)
        if i + EMBED_BATCH_SIZE < total:
            time.sleep(DELAY_BETWEEN_BATCHES / max(rotator.n, 1))
    return all_vectors


# ============================================================
# MAIN
# ============================================================
def main():
    load_dotenv()
    api_keys = load_api_keys()
    if not api_keys:
        print("❌ Lỗi: Không tìm thấy GEMINI_API_KEY hoặc GEMINI_API_KEYS trong file .env")
        sys.exit(1)
    if len(api_keys) > 1:
        print(
            f"🔑 Đang dùng {len(api_keys)} API key, xoay vòng để tránh rate limit.")

    force_rebuild = "--rebuild" in sys.argv
    if force_rebuild:
        print("⚠️  Chế độ REBUILD TOÀN BỘ được bật (--rebuild). Sẽ xoá manifest + index cũ và nhúng lại từ đầu.")
        if os.path.exists(MANIFEST_PATH):
            os.remove(MANIFEST_PATH)

    if not os.path.isdir(DATA_FOLDER):
        print(f"❌ Không tìm thấy thư mục '{DATA_FOLDER}'.")
        sys.exit(1)

    manifest = {} if force_rebuild else load_manifest()

    files_on_disk = {
        f for f in os.listdir(DATA_FOLDER)
        if os.path.isfile(os.path.join(DATA_FOLDER, f))
    }

    # 1) File đã bị xoá khỏi learning_data/ -> cần xoá khỏi index
    removed_files = set(manifest.keys()) - files_on_disk

    # 2) File mới hoặc đã thay đổi nội dung (so hash)
    files_to_process = []
    unchanged_count = 0
    for filename in sorted(files_on_disk):
        path = os.path.join(DATA_FOLDER, filename)
        h = compute_file_hash(path)
        old_entry = manifest.get(filename)
        if old_entry is None or old_entry.get("file_hash") != h:
            files_to_process.append((filename, path, h))
        else:
            unchanged_count += 1

    print(f"📊 Tổng quan: {len(files_on_disk)} file trong '{DATA_FOLDER}' | "
          f"{unchanged_count} không đổi | {len(files_to_process)} mới/thay đổi | "
          f"{len(removed_files)} đã bị xoá")

    if not files_to_process and not removed_files:
        print("✅ Không có gì thay đổi. faiss_index đã cập nhật đầy đủ, không cần chạy lại embedding.")
        return

    rotator = KeyRotator(api_keys)
    # embeddings đại diện (key đầu tiên) chỉ dùng để FAISS biết cách embed câu hỏi lúc query,
    # KHÔNG liên quan tới việc key nào được dùng để tạo vector lúc nhúng dữ liệu (đã xoay vòng riêng).
    embeddings = rotator.current()

    # Nạp index cũ nếu có, để CẬP NHẬT THÊM thay vì xây lại từ đầu
    vectorstore = None
    if os.path.exists(INDEX_FOLDER) and not force_rebuild:
        try:
            vectorstore = FAISS.load_local(
                INDEX_FOLDER, embeddings, allow_dangerous_deserialization=True
            )
            print("📦 Đã nạp faiss_index hiện có — chạy ở chế độ cập nhật (incremental).")
        except Exception as e:
            print(
                f"⚠️  Không nạp được faiss_index cũ ({e}). Sẽ tạo mới từ đầu.")
            vectorstore = None

    # Xoá vector cũ của các file bị xoá / thay đổi (để tránh dữ liệu trùng/lỗi thời)
    ids_to_delete = []
    for filename in removed_files:
        ids_to_delete.extend(manifest[filename].get("chunk_ids", []))
        del manifest[filename]
    for filename, _, _ in files_to_process:
        if filename in manifest:
            ids_to_delete.extend(manifest[filename].get("chunk_ids", []))

    if vectorstore is not None and ids_to_delete:
        try:
            vectorstore.delete(ids=ids_to_delete)
            print(
                f"🗑️  Đã xoá {len(ids_to_delete)} đoạn cũ (thuộc file bị sửa/xoá) khỏi index.")
        except Exception as e:
            print(f"⚠️  Không xoá được vector cũ ({e}). "
                  f"Nếu dữ liệu bị trùng lặp, hãy chạy lại với: python prepare_data.py --rebuild")

    if removed_files:
        save_manifest(manifest)
        if vectorstore is not None:
            vectorstore.save_local(INDEX_FOLDER)

    if not files_to_process:
        print("🎉 Hoàn tất xoá dữ liệu cũ, không có file mới cần nhúng.")
        return

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000, chunk_overlap=200)

    print(f"\n🚀 Bắt đầu xử lý {len(files_to_process)} file mới/thay đổi "
          f"(batch={EMBED_BATCH_SIZE}, nghỉ {DELAY_BETWEEN_BATCHES}s giữa các batch)...\n")

    for filename, path, h in files_to_process:
        print(f"-> Đang trích xuất: {filename}")
        content = extract_text_from_file(path)
        if not content.strip():
            print(f"   [!] Không trích xuất được nội dung, bỏ qua file này.\n")
            continue

        doc = Document(page_content=content, metadata={"source": filename})
        chunks = text_splitter.split_documents([doc])

        chunk_ids = [f"{filename}::{idx}" for idx in range(len(chunks))]
        chunk_texts = [c.page_content for c in chunks]
        metadatas = [c.metadata for c in chunks]

        print(f"   -> Chia thành {len(chunks)} đoạn nhỏ.")

        try:
            vectors = embed_documents_in_batches(rotator, chunk_texts)
        except RuntimeError as e:
            # Hết retry vẫn lỗi -> dừng lại an toàn, giữ nguyên tiến độ đã lưu từ các file trước
            print(f"\n❌ {e}")
            print(
                f"   File '{filename}' và các file sau đó trong danh sách CHƯA được xử lý.")
            sys.exit(1)

        text_embedding_pairs = list(zip(chunk_texts, vectors))

        if vectorstore is None:
            vectorstore = FAISS.from_embeddings(
                text_embedding_pairs, embeddings, metadatas=metadatas, ids=chunk_ids
            )
        else:
            vectorstore.add_embeddings(
                text_embedding_pairs, metadatas=metadatas, ids=chunk_ids
            )

        manifest[filename] = {
            "file_hash": h,
            "chunk_ids": chunk_ids,
            "num_chunks": len(chunk_ids),
        }

        # Checkpoint: lưu ngay sau mỗi file để không mất tiến độ nếu bị lỗi/quota giữa chừng
        vectorstore.save_local(INDEX_FOLDER)
        save_manifest(manifest)
        print(f"   ✅ Đã nhúng và lưu '{filename}' ({len(chunk_ids)} đoạn).\n")

    print("🎉 Hoàn tất! faiss_index đã được cập nhật với dữ liệu mới nhất.")


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n⏹️  Đã dừng theo yêu cầu (Ctrl+C). "
              "Tiến độ tới file gần nhất đã checkpoint vẫn được giữ nguyên trong "
              f"'{INDEX_FOLDER}' và '{MANIFEST_PATH}' — chạy lại script để tiếp tục từ đó.")
        sys.exit(130)
